import sys
import os
import json
from PyQt5.QtWidgets import (QApplication, QWidget, QLabel, QPushButton, QVBoxLayout, QFileDialog, QMessageBox, QComboBox, QLineEdit)
from docx2pdf import convert
from PIL import Image
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from io import BytesIO
from pptx import Presentation

class PdfWatermarker(QWidget):
    def __init__(self):
        super().__init__()

        # Initialize selected files list
        self.selected_files = []
        self.watermark_path = ""
        self.watermark_text = ""

        # Load previous watermark settings if they exist
        self.load_watermark_settings()

        # Set up the UI
        self.init_ui()

    def init_ui(self):
        self.setWindowTitle("File to PDF Converter with Watermark")

        # Create layout and widgets
        layout = QVBoxLayout()

        self.file_label = QLabel("No DOCX or PPTX files selected")
        layout.addWidget(self.file_label)

        choose_file_btn = QPushButton("Choose DOCX or PPTX File(s)")
        choose_file_btn.clicked.connect(self.choose_files)
        layout.addWidget(choose_file_btn)

        self.watermark_type_combo = QComboBox()
        self.watermark_type_combo.addItems(["Image Watermark", "Text Watermark"])
        self.watermark_type_combo.setCurrentText(self.watermark_type)  # Load saved watermark type
        self.watermark_type_combo.currentTextChanged.connect(self.toggle_watermark_input)
        layout.addWidget(self.watermark_type_combo)

        self.watermark_label = QLabel(self.watermark_path if self.watermark_path else "No watermark image selected")
        layout.addWidget(self.watermark_label)

        self.choose_watermark_btn = QPushButton("Choose Watermark Image")
        self.choose_watermark_btn.clicked.connect(self.choose_watermark_file)
        layout.addWidget(self.choose_watermark_btn)

        self.watermark_text_input = QLineEdit()
        self.watermark_text_input.setPlaceholderText("Enter watermark text")
        self.watermark_text_input.setText(self.watermark_text)  # Load saved watermark text
        self.watermark_text_input.setVisible(self.watermark_type == "Text Watermark")
        layout.addWidget(self.watermark_text_input)

        convert_btn = QPushButton("Convert and Apply Watermark")
        convert_btn.clicked.connect(self.convert_and_watermark)
        layout.addWidget(convert_btn)

        self.setLayout(layout)

    def choose_files(self):
        file_paths, _ = QFileDialog.getOpenFileNames(self, "Select DOCX or PPTX Files", "", "Word or PowerPoint files (*.docx;*.pptx)")
        if file_paths:
            self.selected_files = file_paths
            self.file_label.setText(f"{len(file_paths)} file(s) selected")

    def toggle_watermark_input(self, text):
        if text == "Text Watermark":
            self.watermark_text_input.setVisible(True)
            self.choose_watermark_btn.setVisible(False)
        else:
            self.watermark_text_input.setVisible(False)
            self.choose_watermark_btn.setVisible(True)

    def choose_watermark_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Select Watermark Image", "", "Image files (*.png;*.jpg)")
        if file_path:
            self.watermark_path = file_path
            self.watermark_label.setText(file_path)

    def add_image_watermark(self, input_pdf, output_pdf):
        # Open the watermark image using Pillow
        watermark = Image.open(self.watermark_path).convert("RGBA")

        # Apply washout effect (adjust opacity)
        alpha = watermark.split()[3]  # Get the alpha channel
        alpha = alpha.point(lambda p: p * 0.3)  # Reduce opacity to 30%
        watermark.putalpha(alpha)

        # Save the transparent watermark as a temporary image
        temp_watermark_path = "temp_watermark.png"
        watermark.save(temp_watermark_path)

        # Create a BytesIO buffer to store the watermark as a PDF
        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=A4)

        # Set the watermark image in the center (original watermark)
        can.drawImage(temp_watermark_path, 250, 360, width=80, height=100, mask='auto')

        # Add the watermark image in the header on the right side, 2cm below the top
        can.drawImage(temp_watermark_path, 500, 700, width=40, height=50, mask='auto')

        can.save()

        # Move to the beginning of the BytesIO buffer
        packet.seek(0)

        # Read the new PDF with the watermark
        new_pdf = PdfReader(packet)

        # Read the existing PDF
        existing_pdf = PdfReader(open(input_pdf, "rb"))

        # Create a PdfWriter object to write the final output
        output = PdfWriter()

        # Merge the watermark onto each page of the existing PDF
        for i in range(len(existing_pdf.pages)):
            page = existing_pdf.pages[i]
            page.merge_page(new_pdf.pages[0])  # Apply watermark on each page
            output.add_page(page)

        # Write the final watermarked PDF to the specified file
        with open(output_pdf, "wb") as f:
            output.write(f)

        # Clean up the temporary watermark image
        os.remove(temp_watermark_path)

    def add_text_watermark(self, input_pdf, output_pdf, watermark_text):
        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=A4)

        # Set the watermark text at the desired position and size
        can.setFont("Helvetica", 40)
        can.setFillColorRGB(0.5, 0.5, 0.5)  # Light grey color
        can.drawString(250, 400, watermark_text)
        can.save()

        packet.seek(0)
        new_pdf = PdfReader(packet)

        existing_pdf = PdfReader(open(input_pdf, "rb"))
        output = PdfWriter()

        for i in range(len(existing_pdf.pages)):
            page = existing_pdf.pages[i]
            page.merge_page(new_pdf.pages[0])
            output.add_page(page)

        with open(output_pdf, "wb") as f:
            output.write(f)

    def convert_pptx_to_pdf(self, pptx_path, pdf_path):
        ppt = Presentation(pptx_path)
        slides = ppt.slides

        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=A4)

        for slide in slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    can.drawString(100, 700, shape.text)
            can.showPage()

        can.save()

        with open(pdf_path, "wb") as f:
            f.write(packet.getvalue())

    def convert_and_watermark(self):
        if not self.selected_files or (self.watermark_type_combo.currentText() == "Image Watermark" and not self.watermark_path):
            QMessageBox.critical(self, "Error", "Please select files and a watermark option!")
            return

        watermarked_pdf_paths = []

        for file_index, file_path in enumerate(self.selected_files):
            pdf_path = f'{os.path.splitext(file_path)[0]}.pdf'
            watermarked_pdf_path = f'{os.path.splitext(file_path)[0]}_do_not_copy.pdf'

            try:
                print(f"Processing file {file_index + 1}/{len(self.selected_files)}: {file_path}")  # Debug log

                if file_path.endswith(".docx"):
                    # Convert DOCX to PDF
                    try:
                        try:
                            convert(file_path, pdf_path)
                        except Exception as e:
                            convert(file_path, pdf_path)
                    except Exception as e:
                        QMessageBox.critical(self, "Error", f"Failed to convert '{file_path}' to PDF. Error: {str(e)}")
                        continue  # Skip to the next file

                elif file_path.endswith(".pptx"):
                    # Convert PPTX to PDF
                    self.convert_pptx_to_pdf(file_path, pdf_path)

                # Check if the PDF is created before adding a watermark
                if os.path.exists(pdf_path):
                    if self.watermark_type_combo.currentText() == "Image Watermark":
                        self.add_image_watermark(pdf_path, watermarked_pdf_path)
                    else:
                        self.add_text_watermark(pdf_path, watermarked_pdf_path, self.watermark_text_input.text())

                    # Check if watermarked PDF was created successfully
                    if os.path.exists(watermarked_pdf_path):
                        watermarked_pdf_paths.append(watermarked_pdf_path)
                        # Delete the original PDF immediately after watermark is applied
                        os.remove(pdf_path)
                    else:
                        QMessageBox.critical(self, "Error", f"Failed to create watermarked PDF for '{file_path}'")
                else:
                    QMessageBox.critical(self, "Error", f"PDF file not found for '{file_path}' after conversion.")

            except Exception as e:
                pass
                # QMessageBox.critical(self, "Error", f"An error occurred while processing '{file_path}': {str(e)}")
                continue  # Continue to the next file even if one fails

        if watermarked_pdf_paths:
            QMessageBox.information(self, "Success", f"Watermarked PDFs created successfully!\nCreated files: {', '.join(watermarked_pdf_paths)}")

        # Save watermark settings after conversion
        self.save_watermark_settings()

    def save_watermark_settings(self):
        # Save the watermark settings to a JSON file
        settings = {
            "watermark_type": self.watermark_type_combo.currentText(),
            "watermark_path": self.watermark_path,
            "watermark_text": self.watermark_text_input.text()
        }
        with open("watermark_settings.json", "w") as file:
            json.dump(settings, file)

    def load_watermark_settings(self):
        # Load the watermark settings from a JSON file, if it exists
        if os.path.exists("watermark_settings.json"):
            with open("watermark_settings.json", "r") as file:
                settings = json.load(file)
                self.watermark_type = settings.get("watermark_type", "Image Watermark")
                self.watermark_path = settings.get("watermark_path", "")
                self.watermark_text = settings.get("watermark_text", "")
        else:
            self.watermark_type = "Image Watermark"
            self.watermark_path = ""
            self.watermark_text = ""

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = PdfWatermarker()
    window.resize(400, 300)
    window.show()
    sys.exit(app.exec_())
